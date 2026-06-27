import os
from pptx import Presentation
from pptx.util import Inches, Pt

def edit_presentation():
    template_path = r"d:\Projects\ZK-Tester\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
    output_path = r"d:\Projects\ZK-Tester\[Pub] ISRO BAH 2026 _ Idea Submission Template_Filled.pptx"
    
    print(f"Loading PowerPoint template: {template_path}")
    prs = Presentation(template_path)
    
    # ----------------------------------------------------
    # Slide 1: Cover slide
    # ----------------------------------------------------
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "Team Name" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Team Name : Antigravity"
                p.font.size = Pt(28)
                p.font.bold = True
                
                p2 = shape.text_frame.add_paragraph()
                p2.text = "Problem Statement : Satellite Image Super-Resolution & Natural Colorization"
                p2.font.size = Pt(24)
                
                p3 = shape.text_frame.add_paragraph()
                p3.text = "Team Leader Name : [Enter Leader Name]"
                p3.font.size = Pt(24)

    # ----------------------------------------------------
    # Slide 3: Opportunity / USP
    # ----------------------------------------------------
    slide = prs.slides[2]
    # Find the shape that contains "Opportunity" or "USP"
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "Opportunity should be able" in text or "USP" in text:
                shape.text_frame.clear()
                
                p = shape.text_frame.paragraphs[0]
                p.text = "Proposed Solution & Opportunity Mapping"
                p.font.size = Pt(24)
                p.font.bold = True
                p.space_after = Pt(14)
                
                bullets = [
                    ("How different is it from existing ideas?", 0),
                    ("Most colorizers use standard 8-bit RGB/PNG mappings that clip critical sub-meter sensor data. Our pipeline processes raw 16-bit float satellite bands directly in-memory, preserving the full dynamic range of surface reflectance and Kelvin temperatures.", 1),
                    ("How will it solve the problem?", 0),
                    ("Uses a deep cascaded cascade: SRResNet restores high-frequency spatial structures (upscaling IR by 4x), while a bridged UNetColorizer translates thermal signatures into natural visible spectrum RGB colors.", 1),
                    ("USP (Unique Selling Proposition)", 0),
                    ("A fully georeferenced pipeline featuring dynamic Homoscedastic Uncertainty loss weight balancing, Gradient-Domain SSIM edge enforcement, and overlap Cosine Seam Feathering for artifact-free mosaic stitching.", 1)
                ]
                
                for b_text, level in bullets:
                    p = shape.text_frame.add_paragraph()
                    p.text = b_text
                    p.level = level
                    p.font.size = Pt(18 - 2 * level)
                    if level == 0:
                        p.font.bold = True
                        p.space_before = Pt(8)
                    p.space_after = Pt(4)

    # ----------------------------------------------------
    # Slide 4: Features
    # ----------------------------------------------------
    slide = prs.slides[3]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "List of features offered" in text:
                shape.text_frame.clear()
                
                p = shape.text_frame.paragraphs[0]
                p.text = "Core Features of the Solution"
                p.font.size = Pt(24)
                p.font.bold = True
                p.space_after = Pt(14)
                
                features = [
                    ("1. Raw 16-Bit float Preservation:", 0),
                    ("Retains raw satellite sensor values without visual compression, quantization, or loss of thermal detail.", 1),
                    ("2. Sub-Pixel Upsampling Cascade:", 0),
                    ("Achieves high-quality 4x spatial super-resolution using PixelShuffle upsampling steerable by structural skip-connections.", 1),
                    ("3. Boundary Edge Enforcement:", 0),
                    ("Minimizes SSIM over horizontal and vertical Sobel gradient maps to ensure razor-sharp borders on color boundaries.", 1),
                    ("4. Overlap Cosine Blending:", 0),
                    ("Uses a 2D cosine windowing function to blend overlapping tile seams, eliminating checkerboard lines.", 1),
                    ("5. Dynamic Multitask Loss Weighting:", 0),
                    ("Auto-balances L1, gradient, and semantic consistency losses dynamically during training using learnable log-variances.", 1)
                ]
                
                for b_text, level in features:
                    p = shape.text_frame.add_paragraph()
                    p.text = b_text
                    p.level = level
                    p.font.size = Pt(18 - 2 * level)
                    if level == 0:
                        p.font.bold = True
                        p.space_before = Pt(8)
                    p.space_after = Pt(4)

    # ----------------------------------------------------
    # Slide 5: Process Flow Diagram
    # ----------------------------------------------------
    slide = prs.slides[4]
    # Delete template text placeholder
    to_delete = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "Process flow" in text or "Add a flow" in text:
                to_delete.append(shape)
                
    for shape in to_delete:
        # We can remove shape
        sp = shape._element
        sp.getparent().remove(sp)
        
    # Insert pipeline diagram PNG
    pipeline_img = r"d:\Projects\ZK-Tester\IR-Image-Colouring\outputs\diag_pipeline.png"
    if os.path.exists(pipeline_img):
        print("Inserting pipeline diagram on Slide 5...")
        # Center of slide: width=13.33, height=7.5
        left = Inches(1.5)
        top = Inches(1.2)
        width = Inches(10.33)
        height = Inches(5.5)
        slide.shapes.add_picture(pipeline_img, left, top, width=width, height=height)

    # ----------------------------------------------------
    # Slide 7: Architecture Diagram
    # ----------------------------------------------------
    slide = prs.slides[6]
    to_delete = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "Architecture diagram" in text:
                to_delete.append(shape)
                
    for shape in to_delete:
        sp = shape._element
        sp.getparent().remove(sp)
        
    # Insert cascade diagram PNG
    cascade_img = r"d:\Projects\ZK-Tester\IR-Image-Colouring\outputs\diag_cascade.png"
    if os.path.exists(cascade_img):
        print("Inserting cascade diagram on Slide 7...")
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(11.33)
        height = Inches(5.0)
        slide.shapes.add_picture(cascade_img, left, top, width=width, height=height)

    # ----------------------------------------------------
    # Slide 8: Technologies Used
    # ----------------------------------------------------
    slide = prs.slides[7]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs)
            if "Technologies to be used" in text:
                shape.text_frame.clear()
                
                p = shape.text_frame.paragraphs[0]
                p.text = "Technologies & Frameworks"
                p.font.size = Pt(24)
                p.font.bold = True
                p.space_after = Pt(14)
                
                techs = [
                    ("Deep Learning Core Framework:", 0),
                    ("PyTorch & PyTorch Lightning (Mixed-precision 16-bit GPU acceleration).", 1),
                    ("Geospatial & Remote Sensing Libraries:", 0),
                    ("Rasterio & GDAL (For co-registration, warping, windowed tiling, and georeference preservation).", 1),
                    ("Computer Vision & Objective Losses:", 0),
                    ("Kornia, OpenCV, Albumentations (Gradient-domain Sobel operations, image metrics).", 1),
                    ("Inference Compiler Optimization:", 0),
                    ("Nvidia TensorRT Hook placeholder (Compilation to FP16 execution graph).", 1),
                    ("Experiment Logging & Tracking:", 0),
                    ("Weights & Biases (W&B) / TensorBoard.", 1)
                ]
                
                for b_text, level in techs:
                    p = shape.text_frame.add_paragraph()
                    p.text = b_text
                    p.level = level
                    p.font.size = Pt(18 - 2 * level)
                    if level == 0:
                        p.font.bold = True
                        p.space_before = Pt(8)
                    p.space_after = Pt(4)
                    
    print(f"Saving edited PowerPoint presentation: {output_path}")
    prs.save(output_path)
    print("Done!")

if __name__ == "__main__":
    edit_presentation()
