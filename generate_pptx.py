import os
import base64
import time
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Mermaid diagram syntaxes
DIAGRAM_PIPELINE = """graph TD
    B4[Visible Red Band B4 - 30m]
    B3[Visible Green Band B3 - 30m]
    B2[Visible Blue Band B2 - 30m]
    B5[Infrared Band B5 - 30m]
    B10[Thermal Band B10 - 100m]

    RS[Warp Reprojection <br> Resampling.cubic_spline]
    Stack[Stack RGB Bands]
    Tiler[Windowed Tiling <br> 512x512 with 64 overlap]
    TIFF[Write 16-bit Float GeoTIFFs <br> Preserving CRS + Geotransform]

    B10 --> RS
    B5 --> RS
    B4 --> Stack
    B3 --> Stack
    B2 --> Stack
    RS --> Tiler
    Stack --> Tiler
    Tiler --> TIFF

    Norm[Local Z-Score Normalization]
    SR[SRResNet <br> 4x Super-Resolution]
    Color[UNetColorizer <br> IR to RGB Translation]

    TIFF -- "Raw IR Tile" --> Norm
    Norm -- "Normalized IR" --> SR
    SR -- "Upscaled Structure" --> Color
    SR -- "Low-Res Bridge Features" --> Color
    Color -- "RGB Prediction" --> Stitch[Overlap Cosine Blending]

    ScaleGeo[Scale Geotransform by 1/4]
    OutTIFF[Georeferenced RGB GeoTIFF <br> float32]

    Stitch --> ScaleGeo
    ScaleGeo --> OutTIFF
"""

DIAGRAM_CASCADE = """graph LR
    subgraph SRResNet Module
        InIR["Low-Res IR Input <br> (B, 2, 128, 128)"] --> Conv1["Conv 9x9 + PReLU"]
        Conv1 --> ResBlocks["8 Residual Blocks"]
        ResBlocks --> Conv2["Conv 3x3 + BatchNorm"]
        
        ResBlocks -- "Bridge Features <br> (B, 64, 128, 128)" --> ConcatPoint
        
        Conv2 --> PS["PixelShuffle Upsampler <br> (4x Spatial Scaling)"]
        PS --> ConvOut["Conv 9x9 (Output)"]
        ConvOut --> HighResIR["Upscaled IR prior <br> (B, 2, 512, 512)"]
    end

    subgraph UNetColorizer Module
        HighResIR --> Enc1["Encoder Level 1 <br> (512 -> 256)"]
        Enc1 --> Enc2["Encoder Level 2 <br> (256 -> 128)"]
        
        Enc2 --> ConcatPoint["Channel Concatenation <br> (Total: 128 channels)"]
        ConcatPoint --> Enc3["Encoder Level 3 <br> (128 -> 64)"]
        Enc3 --> Enc4["Encoder Level 4 <br> (64 -> 32)"]
        
        Enc4 --> Bottleneck["Bottleneck Layers <br> (Channel: 512)"]
        
        Bottleneck --> Dec4["Decoder Level 4 <br> (32 -> 64)"]
        Dec4 -- "Skip" --> Enc4
        Dec4 --> Dec3["Decoder Level 3 <br> (64 -> 128)"]
        Dec3 -- "Skip" --> Enc3
        Dec3 --> Dec2["Decoder Level 2 <br> (128 -> 256)"]
        Dec2 -- "Skip" --> Enc2
        Dec2 --> Dec1["Decoder Level 1 <br> (256 -> 512)"]
        Dec1 -- "Skip" --> Enc1
        
        Dec1 --> Final["Final Conv 1x1"]
        Final --> RGB["Visible RGB Output <br> (B, 3, 512, 512)"]
    end
"""

DIAGRAM_LOSS = """graph TD
    Pred["Predicted RGB"] --> L1[L1 Pixel Reconstruction Loss]
    GT["Ground Truth RGB"] --> L1

    Pred --> Grad[Gradient-Domain SSIM Loss]
    GT --> Grad

    Pred --> Sem[Semantic Consistency Loss]
    GT --> Sem
    Frozen["Frozen SegFormer Backbone"] --> Sem

    L1 --> L_total["Weighted Total Loss Optimizer"]
    Grad --> L_total
    Sem --> L_total

    subgraph Uncertainty Weighting Engine
        LV1["Learnable Log-Var s_1"] -- "exp(-s_1)" --> L_total
        LV2["Learnable Log-Var s_2"] -- "exp(-s_2)" --> L_total
        LV3["Learnable Log-Var s_3"] -- "exp(-s_3)" --> L_total
        
        LV1 -- "+ s_1" --> L_total
        LV2 -- "+ s_2" --> L_total
        LV3 -- "+ s_3" --> L_total
    end
"""

def download_diagram(diagram_syntax, filename):
    encoded = base64.b64encode(diagram_syntax.encode('utf-8')).decode('utf-8')
    url = f"https://mermaid.ink/img/{encoded}"
    
    # Try downloading with retries
    for attempt in range(3):
        print(f"Rendering and downloading {filename} via mermaid.ink (Attempt {attempt+1}/3)...")
        try:
            response = requests.get(url, timeout=60)
            if response.status_code == 200:
                os.makedirs("outputs", exist_ok=True)
                path = f"outputs/{filename}"
                with open(path, "wb") as f:
                    f.write(response.content)
                print(f"Downloaded successfully: {path}")
                return path
            else:
                print(f"Mermaid.ink returned status code: {response.status_code}")
        except Exception as e:
            print(f"Error on attempt {attempt+1}: {e}")
        time.sleep(2) # wait before retrying
    return None

def create_presentation():
    # Download diagrams
    pipeline_img = download_diagram(DIAGRAM_PIPELINE, "diag_pipeline.png")
    cascade_img = download_diagram(DIAGRAM_CASCADE, "diag_cascade.png")
    loss_img = download_diagram(DIAGRAM_LOSS, "diag_loss.png")
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    BG_COLOR = RGBColor(20, 24, 30)      # Deep dark slate
    TITLE_COLOR = RGBColor(0, 200, 200)   # Cyan / Mint
    TEXT_COLOR = RGBColor(230, 235, 240)  # Off-white
    MUTED_COLOR = RGBColor(170, 180, 190) # Muted grey
    
    def set_slide_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_slide_header(slide, title_text):
        set_slide_background(slide)
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Underline accent shape
        rect = slide.shapes.add_shape(1, Inches(0.75), Inches(1.3), Inches(3.0), Inches(0.04)) # 1 = RECTANGLE
        rect.fill.solid()
        rect.fill.fore_color.rgb = TITLE_COLOR
        rect.line.fill.background()
        
    def add_bullet_content(slide, bullets, left=Inches(0.75), top=Inches(1.8), width=Inches(11.83), height=Inches(5.0), size=18):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for i, item in enumerate(bullets):
            if isinstance(item, tuple):
                b_text, level = item
            else:
                b_text, level = item, 0
                
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                
            p.text = b_text
            p.level = level
            p.font.name = 'Segoe UI'
            p.font.size = Pt(size - 2 * level)
            if level == 0:
                p.font.color.rgb = TEXT_COLOR
                p.font.bold = True
                p.space_before = Pt(12)
            else:
                p.font.color.rgb = MUTED_COLOR
                p.space_before = Pt(4)
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Custom layout)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1)
    
    # Large Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.text = "Deep Geospatial Cascade Model"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TITLE_COLOR
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Infrared Super-Resolution & Colorization (v2.0)"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(32)
    p2.font.color.rgb = TEXT_COLOR
    p2.space_after = Pt(30)
    
    p3 = tf.add_paragraph()
    p3.text = "Translating Landsat 9 Multi-spectral Sensor Data to Seamless GIS-Aligned RGB"
    p3.font.name = 'Segoe UI'
    p3.font.size = Pt(18)
    p3.font.color.rgb = MUTED_COLOR
    
    # Accent footer
    line = slide1.shapes.add_shape(1, Inches(1.0), Inches(5.2), Inches(4.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TITLE_COLOR
    line.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 2: The Core Challenge
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "The Resolution and Interpretation Bottleneck")
    bullets2 = [
        ("The Spectral Interpretation Gap", 0),
        ("Infrared and Thermal bands (NIR, TIRS) contain crucial physical signatures (soil moisture, heat, vegetation), but they are grayscale and unintuitive for human verification.", 1),
        ("The Spatial Resolution Discrepancy", 0),
        ("Thermal bands have lower native resolution (100m) compared to visible bands (30m). Scaling them directly creates pixelated artifacts.", 1),
        ("The Spatial Preservation Requirement", 0),
        ("Any translation must completely preserve physical coordinates, Coordinate Reference Systems (CRS), and geotransform grids to remain GIS-compatible.", 1)
    ]
    add_bullet_content(slide2, bullets2)

    # ----------------------------------------------------
    # SLIDE 3: System Pipeline Architecture (with Diagram)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "System Pipeline Architecture")
    bullets3 = [
        ("1. Preprocessing Stage (Geospatial Ingest)", 0),
        ("Warp-aligns bands onto a unified visible grid using Resampling.cubic_spline and crops them into 512x512 tiles.", 1),
        ("2. Deep Cascade Inference (Neural Translation)", 0),
        ("Runs SRResNet on 128x128 decimated inputs, passing upscaled structure and intermediate bridge features to the UNetColorizer.", 1),
        ("3. Stitching & Metadata Restoration", 0),
        ("Applies Overlap Cosine Seam Feathering to blend borders, scales the affine geotransform, and writes out a clean 16-bit float GeoTIFF.", 1)
    ]
    
    if pipeline_img:
        # Split layout: text on left, image on right
        add_bullet_content(slide3, bullets3, left=Inches(0.75), top=Inches(1.8), width=Inches(5.8), height=Inches(5.0), size=17)
        # Position image nicely on right half
        slide3.shapes.add_picture(pipeline_img, Inches(6.8), Inches(1.5), width=Inches(5.8), height=Inches(5.3))
    else:
        # Full width layout fallback
        add_bullet_content(slide3, bullets3)

    # ----------------------------------------------------
    # SLIDE 4: Precision Geospatial Preprocessing
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "Precision Geospatial Ingestion")
    bullets4 = [
        ("Preserving the 16-Bit High Dynamic Range (HDR)", 0),
        ("Warped tiles are written directly as float32 arrays. We bypass lossy 8-bit PNG conversions that clip or compress the sensor data.", 1),
        ("Warp Co-Registration", 0),
        ("Thermal bands (B10) are re-projected onto the Visible reference (B4) grid dynamically using Rasterio warp engines.", 1),
        ("In-Memory Z-Score Normalization", 0),
        ("Inputs are adapted to the model dynamically using per-tile statistics: (x - mean) / (std + epsilon). This maps raw values (1000s) to active zero-centered distributions.", 1)
    ]
    add_bullet_content(slide4, bullets4)

    # ----------------------------------------------------
    # SLIDE 5: SRResNet: Structure Recovery
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "SRResNet: Restoring Spatial Structure")
    bullets5 = [
        ("Core Functionality", 0),
        ("Takes low-res IR bands (decimated by 4x to 128x128) and upscales them to reconstruct sharp boundary borders.", 1),
        ("Deep Residual Cascade", 0),
        ("Uses 8 Residual Blocks (Conv2d -> BatchNorm2d -> PReLU -> Conv2d -> BatchNorm2d) with identity shortcuts to stabilize gradient flow.", 1),
        ("Sub-Pixel Upsampling & Features Bridge", 0),
        ("Upscaling is handled by sub-pixel convolutions (PixelShuffle layers). It hooks intermediate low-resolution feature maps (64 channels) to pass to the colorizer.", 1)
    ]
    add_bullet_content(slide5, bullets5)

    # ----------------------------------------------------
    # SLIDE 6: UNetColorizer: Bridged Translation
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "UNetColorizer: Feature-Bridged Translation")
    bullets6 = [
        ("Bridged Cross-Attention / Skip Concat", 0),
        ("Directly injects the 64-channel bridge features from SRResNet into the U-Net bottleneck level (128x128 resolution), guiding colorization with structural cues.", 1),
        ("Encoder-Decoder Architecture", 0),
        ("Uses a 4-level encoder (convolution + pooling down to 32x32 bottleneck) and 4-level decoder with skip-connections to reconstruct spatial context.", 1),
        ("Natural Color Mapping", 0),
        ("Translates the structural IR prior into natural visible red, green, and blue (RGB) spectrum bands.", 1)
    ]
    add_bullet_content(slide6, bullets6)

    # ----------------------------------------------------
    # SLIDE 7: Model Architecture Diagram (Dedicated full-slide)
    # ----------------------------------------------------
    if cascade_img:
        slide7 = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide7, "Model Architecture Diagram")
        # Add full-screen cascade diagram
        slide7.shapes.add_picture(cascade_img, Inches(1.0), Inches(1.7), width=Inches(11.33), height=Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 8: Loss Formulations & Objectives
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "Multi-Domain Loss Formulations")
    bullets8 = [
        ("L1 Reconstruction Loss", 0),
        ("Measures pixel-to-pixel absolute errors between the predicted RGB and ground truth Visible channels to establish base color structures.", 1),
        ("Gradient-Domain SSIM Loss (Edge Detail)", 0),
        ("Extracts horizontal and vertical gradients (Sobel filters) of prediction and ground truth. Minimizes (1 - SSIM) of gradient maps to enforce sharp boundaries.", 1),
        ("Semantic Consistency Loss (Guardrail)", 0),
        ("Evaluates land-cover class segmentations on predicted vs ground-truth RGB using a frozen segmenter, penalizing classification shift (hallucination).", 1)
    ]
    add_bullet_content(slide8, bullets8)

    # ----------------------------------------------------
    # SLIDE 9: Multi-Task Loss Auto-Tuning (with Diagram)
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide9, "Dynamic Loss Weight Balancing")
    bullets9 = [
        ("The Challenge: Optimizer Conflicts", 0),
        ("Statically scaling L1, Gradient, and Semantic weights causes optimizer conflicts, where one domain dominates and destabilizes training.", 1),
        ("The Solution: Uncertainty Auto-Tuning", 0),
        ("We parameterize task uncertainty as learnable log-variances (log_var_i) optimized in parallel with model weights.", 1),
        ("Weight Optimization Formula", 0),
        ("Loss = sum_i [ exp(-log_var_i) * Loss_i + 0.5 * log_var_i ]. The system penalizes large weights while dividing each loss dynamically to stabilize convergence.", 1)
    ]
    
    if loss_img:
        # Split layout
        add_bullet_content(slide9, bullets9, left=Inches(0.75), top=Inches(1.8), width=Inches(5.8), height=Inches(5.0), size=17)
        # Position image nicely on right half
        slide9.shapes.add_picture(loss_img, Inches(7.0), Inches(1.6), width=Inches(5.5), height=Inches(5.1))
    else:
        add_bullet_content(slide9, bullets9)

    # ----------------------------------------------------
    # SLIDE 10: Geospatial Blending & Stitching
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide10, "Seamless Stitching & GIS Compatibility")
    bullets10 = [
        ("Overlapping Tile Grids", 0),
        ("Inference is computed on sub-windows with 64-pixel overlaps to accommodate arbitrarily large satellite scenes.", 1),
        ("2D Cosine Seam Blending", 0),
        ("Applies a 2D cosine fade window mask on tile edges: W(t) = 0.5 - 0.5 * cos(pi * t). Smoothly interpolates overlap boundaries, eliminating visual checkerboard seams.", 1),
        ("Geotransform Grid Scaling", 0),
        ("The source geotransform scale is multiplied by (1/scale_factor) to account for upscaling. Ensures output pixels align perfectly when opened in QGIS / ArcGIS.", 1)
    ]
    add_bullet_content(slide10, bullets10)

    # ----------------------------------------------------
    # SLIDE 11: Experimental Results
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide11, "Experimental Results & Validation")
    bullets11 = [
        ("Dataset Parameters", 0),
        ("361 uniform tiles extracted from Landsat 9 scene LC09_L2SP_011001_20260624_20260625_02_T1.", 1),
        ("Deep Cascade Metrics (10-Epoch checkpoint)", 0),
        ("Reconstruction Quality: PSNR of 18.00 dB, SSIM of 0.2731 (structural details converging).", 1),
        ("Semantic Downstream agreement: 88.78% classification pixel agreement.", 1),
        ("Inference Latency: 760.85 ms per tile on NVIDIA GeForce RTX 4050 GPU (5.7M parameters).", 1)
    ]
    add_bullet_content(slide11, bullets11)

    # ----------------------------------------------------
    # SLIDE 12: Summary & Next Steps
    # ----------------------------------------------------
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide12, "Summary & Next Steps")
    bullets12 = [
        ("Completed Scaffold", 0),
        ("Fully operational preprocessing, training, inference, and evaluation pipeline using real deep models with 100% test coverage.", 1),
        ("Next Steps", 0),
        ("1. Complete the 100-epoch training run (~53 mins total) to resolve fine-grained details and realistic colors.", 1),
        ("2. Scale up by importing multiple scenes into raw_dir (preparation and data loaders handle multiple scenes automatically).", 1)
    ]
    add_bullet_content(slide12, bullets12)

    # Save presentation
    os.makedirs("outputs", exist_ok=True)
    out_path = "outputs/project_presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    create_presentation()
